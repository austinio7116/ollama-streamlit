import streamlit as st
import os
import json
from datetime import datetime
from llama_index.core import (
    VectorStoreIndex,
    SimpleDirectoryReader,
    StorageContext,
    load_index_from_storage,
)
from llama_index.llms.ollama import Ollama
from llama_index.embeddings.ollama import OllamaEmbedding
from llama_index.core import Document, Settings
import PyPDF2
import docx
import pptx
import pandas as pd

# ---- Configuration ---- #
DATA_DIR = "./data"
STORAGE_DIR = "./storage"
HISTORY_FILE = "./conversation_history.jsonl"
MODEL = "gemma3"
EMBED_MODEL = "mxbai-embed-large"

# ---- Initialize Ollama LLM and Embeddings ---- #
@st.cache_resource(show_spinner=False)
def init_llm():
    return Ollama(
        model=MODEL,
        request_timeout=180.0,
        stop=["<end_of_turn>"],
        temperature=0.1
    )

@st.cache_resource(show_spinner=False)
def init_embed_model():
    return OllamaEmbedding(model_name=EMBED_MODEL, base_url="http://localhost:11434")

llm = init_llm()
embed_model = init_embed_model()

# ---- Indexing Setup ---- #
def build_or_load_index():
    Settings.embed_model = embed_model
    if not os.path.exists(STORAGE_DIR):
        os.makedirs(STORAGE_DIR, exist_ok=True)
        documents = SimpleDirectoryReader(
            input_dir=DATA_DIR, recursive=True, filename_as_id=True
        ).load_data()
        index = VectorStoreIndex.from_documents(documents)
        index.storage_context.persist(persist_dir=STORAGE_DIR)
    else:
        storage_context = StorageContext.from_defaults(persist_dir=STORAGE_DIR)
        index = load_index_from_storage(storage_context)
    return index

index = build_or_load_index()

# ---- Conversation History ---- #
def load_conversation():
    if os.path.exists(HISTORY_FILE):
        with open(HISTORY_FILE, 'r') as f:
            return [json.loads(line) for line in f]
    return []

def save_conversation(query, response):
    entry = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "query": query,
        "response": response
    }
    with open(HISTORY_FILE, 'a') as f:
        f.write(json.dumps(entry) + "\n")

conversation = load_conversation()

# ---- Hybrid Query Function ---- #
def hybrid_query(prompt, index, llm, conversation, similarity_top_k=3, similarity_cutoff=0.5):
    retriever = index.as_retriever(similarity_top_k=similarity_top_k)
    nodes = retriever.retrieve(prompt)
    print(nodes)

    relevant_nodes = [node for node in nodes if node.score >= similarity_cutoff]
    context = "\n\n".join(node.get_content() for node in relevant_nodes) if relevant_nodes else ""

    conversation_history = ""
    for entry in conversation[-5:]:
        conversation_history += f"<start_of_turn>user\n{entry['query']}<end_of_turn>\n"
        conversation_history += f"<start_of_turn>assistant\n{entry['response']}<end_of_turn>\n"

    conversation_history += f"<start_of_turn>user\n{prompt}<end_of_turn>\n"

    if context:
        conversation_history += f"<start_of_turn>context\n{context}<end_of_turn>\n"

    conversation_history += "<start_of_turn>assistant\n"
    print(conversation_history)
    return llm.stream_complete(conversation_history)

# ---- Streamlit UI ---- #
st.set_page_config(page_title="Gemma3 Multimodal Chat", layout="wide")
st.title("🧠 ChipGPT Multimodal Assistant")

with st.sidebar:
    st.header("📚 Indexed Files")
    for file in os.listdir(DATA_DIR):
        st.write(f"• {file}")

    if st.button("Clear Conversation History"):
        if os.path.exists(HISTORY_FILE):
            os.remove(HISTORY_FILE)
        st.session_state.clear()
        st.rerun()

st.divider()

# Display previous conversation
for entry in conversation[-10:]:
    with st.chat_message("user"):
        st.markdown(entry['query'])
    with st.chat_message("assistant"):
        st.markdown(entry['response'])

# ---- User Input and Optional File Upload ---- #
user_input = st.chat_input("Ask ChipGPT anything:")
uploaded_file = st.file_uploader(
    "Upload a file for additional context (optional):", 
    type=["txt", "pdf", "md", "docx", "pptx", "csv"]
)

def extract_file_content(uploaded_file):
    if uploaded_file.type == "application/pdf":
        import PyPDF2
        reader = PyPDF2.PdfReader(uploaded_file)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    
    elif uploaded_file.type == "text/plain" or uploaded_file.type == "text/markdown":
        return uploaded_file.read().decode("utf-8")
    
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx
        doc = docx.Document(uploaded_file)
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)

    elif uploaded_file.type == "text/csv":
        import pandas as pd
        df = pd.read_csv(uploaded_file)
        return df.to_string(index=False)

    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        from pptx import Presentation
        ppt = Presentation(uploaded_file)
        content = ""
        for slide in ppt.slides:
            content = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            content = "\n".join(content)
        return content

    else:
        # For TXT, MD, or other plain text formats
        return uploaded_file.read().decode("utf-8")

if user_input:
    combined_prompt = user_input

    if uploaded_file:
        file_content = extract_file_content(uploaded_file)
        uploaded_doc = Document(text=file_content)
        index.insert(uploaded_doc)
        combined_prompt += f"\n\nFile Content:\n{file_content}"

    with st.chat_message("user"):
        st.markdown(user_input)
        if uploaded_file:
            st.markdown("📄 Uploaded file included.")

    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        full_response = ""

        streaming_response = hybrid_query(combined_prompt, index, llm, conversation)
        for token in streaming_response:
            full_response += token.delta or ""
            message_placeholder.markdown(full_response)

    save_conversation(user_input, full_response)